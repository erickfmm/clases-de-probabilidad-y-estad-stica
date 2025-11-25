"""
Script para generar presentaciones PowerPoint desde archivos YAML.

Este script lee archivos YAML con contenido de temas de Probabilidad y Estadística,
y genera presentaciones en formato PowerPoint (.pptx).
"""

import yaml
from pathlib import Path
import argparse
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import io


# Colores del tema (basados en la plantilla LaTeX)
COLORES = {
    'primario': RGBColor(41, 128, 185),      # Azul vibrante
    'secundario': RGBColor(231, 76, 60),     # Rojo coral
    'acento': RGBColor(46, 204, 113),        # Verde esmeralda
    'advertencia': RGBColor(241, 196, 15),   # Amarillo dorado
    'morado': RGBColor(155, 89, 182),        # Morado amigable
    'naranja': RGBColor(230, 126, 34),       # Naranja cálido
    'texto': RGBColor(44, 62, 80),           # Gris oscuro para texto
    'fondo_claro': RGBColor(236, 240, 241),  # Gris claro
}


def cargar_yaml(archivo_yaml):
    """
    Carga un archivo YAML y retorna su contenido.
    
    Args:
        archivo_yaml (str): Ruta al archivo YAML
        
    Returns:
        dict: Contenido del archivo YAML
    """
    with open(archivo_yaml, 'r', encoding='utf-8') as f:
        return yaml.safe_load(f)


def crear_presentacion_base():
    """
    Crea una presentación PowerPoint base con configuración inicial.
    
    Returns:
        Presentation: Objeto de presentación de python-pptx
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs


def agregar_portada(prs, tema, subtitulo=""):
    """
    Agrega una diapositiva de portada.
    
    Args:
        prs: Objeto Presentation
        tema (str): Título del tema
        subtitulo (str): Subtítulo opcional
    """
    slide_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo de color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORES['primario']
    
    # Título
    titulo_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    titulo_frame = titulo_box.text_frame
    titulo_frame.word_wrap = True
    p = titulo_frame.paragraphs[0]
    p.text = tema
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    if subtitulo:
        subtitulo_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(0.8)
        )
        subtitulo_frame = subtitulo_box.text_frame
        p = subtitulo_frame.paragraphs[0]
        p.text = subtitulo
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER


def agregar_texto_simple(text_frame, texto, nivel=0, color=None):
    """
    Agrega texto simple al frame de texto.
    
    Args:
        text_frame: Frame de texto donde agregar
        texto (str): Texto a agregar
        nivel (int): Nivel de indentación (0, 1, 2...)
        color (RGBColor): Color del texto
    """
    p = text_frame.add_paragraph()
    p.text = texto
    p.level = nivel
    p.font.size = Pt(18)
    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = COLORES['texto']


def agregar_bloque_especial(text_frame, tipo, texto):
    """
    Agrega un bloque especial (nota, ejemplo, problema, etc.).
    
    Args:
        text_frame: Frame de texto donde agregar
        tipo (str): Tipo de bloque (nota, ejemplo, problema, etc.)
        texto (str): Texto del bloque
    """
    p = text_frame.add_paragraph()
    
    # Iconos y colores según tipo
    iconos_colores = {
        'nota': ('💡 ', COLORES['advertencia']),
        'ejemplo': ('📝 ', COLORES['acento']),
        'problema': ('❓ ', COLORES['secundario']),
        'formula': ('📐 ', COLORES['morado']),
        'calculo': ('🔢 ', COLORES['primario']),
    }
    
    icono, color = iconos_colores.get(tipo, ('• ', COLORES['texto']))
    
    # Agregar icono
    run = p.add_run()
    run.text = icono
    run.font.size = Pt(16)
    
    # Agregar texto
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(16)
    run.font.italic = True
    run.font.color.rgb = color


def agregar_tabla(slide, encabezados, filas, left=1, top=2, width=8, height=4):
    """
    Agrega una tabla a la diapositiva.
    
    Args:
        slide: Diapositiva donde agregar
        encabezados (list): Lista de encabezados
        filas (list): Lista de listas con las filas
        left, top, width, height: Posición y tamaño en pulgadas
    """
    rows = len(filas) + 1  # +1 para encabezados
    cols = len(encabezados)
    
    table = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table
    
    # Encabezados
    for col_idx, encabezado in enumerate(encabezados):
        cell = table.rows[0].cells[col_idx]
        cell.text = str(encabezado)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORES['primario']
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Filas de datos
    for row_idx, fila in enumerate(filas, start=1):
        for col_idx, valor in enumerate(fila):
            cell = table.rows[row_idx].cells[col_idx]
            cell.text = str(valor)
            
            # Alternar colores de fondo
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORES['fondo_claro']
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER


def agregar_grafico_barras(slide, categorias, valores, etiqueta_x="", etiqueta_y="", titulo_serie="Serie"):
    """
    Agrega un gráfico de barras.
    
    Args:
        slide: Diapositiva donde agregar
        categorias (list): Categorías para el eje X
        valores (list): Valores para cada categoría
        etiqueta_x (str): Etiqueta del eje X
        etiqueta_y (str): Etiqueta del eje Y
        titulo_serie (str): Nombre de la serie
    """
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    chart_data.add_series(titulo_serie, valores)
    
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    
    # Personalizar colores
    series = chart.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = COLORES['primario']


def agregar_grafico_lineas(slide, datos_x, datos_y, etiqueta_x="", etiqueta_y="", titulo_serie="Serie"):
    """
    Agrega un gráfico de líneas.
    
    Args:
        slide: Diapositiva donde agregar
        datos_x (list): Valores del eje X
        datos_y (list): Valores del eje Y
        etiqueta_x (str): Etiqueta del eje X
        etiqueta_y (str): Etiqueta del eje Y
        titulo_serie (str): Nombre de la serie
    """
    chart_data = CategoryChartData()
    chart_data.categories = [str(x) for x in datos_x]
    chart_data.add_series(titulo_serie, datos_y)
    
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    
    # Personalizar colores
    series = chart.series[0]
    line = series.format.line
    line.color.rgb = COLORES['secundario']
    line.width = Pt(3)


def agregar_grafico_circular(slide, etiquetas, valores):
    """
    Agrega un gráfico circular (pie chart).
    
    Args:
        slide: Diapositiva donde agregar
        etiquetas (list): Etiquetas para cada porción
        valores (list): Valores para cada porción
    """
    chart_data = CategoryChartData()
    chart_data.categories = etiquetas
    chart_data.add_series('Serie 1', valores)
    
    x, y, cx, cy = Inches(2), Inches(1.5), Inches(6), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False


def procesar_contenido(slide, contenido_items):
    """
    Procesa el contenido de una diapositiva.
    
    Args:
        slide: Diapositiva donde agregar contenido
        contenido_items (list): Lista de items de contenido
    """
    # Separar contenido en texto y gráficos
    tiene_grafico = any(
        isinstance(item, dict) and item.get('tipo') in ['grafico_barras', 'grafico_lineas', 'grafico_circular']
        for item in contenido_items
    )
    
    tiene_tabla = any(
        isinstance(item, dict) and item.get('tipo') == 'tabla'
        for item in contenido_items
    )
    
    # Si hay gráfico o tabla, dividir el espacio
    if tiene_grafico or tiene_tabla:
        # Procesar cada item
        for item in contenido_items:
            if isinstance(item, dict):
                tipo = item.get('tipo', '')
                
                if tipo == 'grafico_barras':
                    agregar_grafico_barras(
                        slide,
                        item.get('categorias', []),
                        item.get('valores', []),
                        item.get('etiqueta_x', ''),
                        item.get('etiqueta_y', ''),
                        item.get('titulo_serie', 'Serie')
                    )
                elif tipo == 'grafico_lineas':
                    agregar_grafico_lineas(
                        slide,
                        item.get('datos_x', []),
                        item.get('datos_y', []),
                        item.get('etiqueta_x', ''),
                        item.get('etiqueta_y', ''),
                        item.get('titulo_serie', 'Serie')
                    )
                elif tipo == 'grafico_circular':
                    agregar_grafico_circular(
                        slide,
                        item.get('etiquetas', []),
                        item.get('valores', [])
                    )
                elif tipo == 'tabla':
                    agregar_tabla(
                        slide,
                        item.get('encabezados', []),
                        item.get('filas', [])
                    )
    else:
        # Solo texto - usar todo el espacio
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(9), Inches(5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Procesar items de contenido
        for item in contenido_items:
            if isinstance(item, str):
                # Texto simple
                agregar_texto_simple(text_frame, item)
            elif isinstance(item, dict):
                tipo = item.get('tipo', '')
                
                if tipo in ['nota', 'ejemplo', 'problema', 'formula', 'calculo']:
                    agregar_bloque_especial(text_frame, tipo, item.get('texto', ''))
                elif tipo == 'componentes':
                    for comp in item.get('lista', []):
                        agregar_texto_simple(text_frame, comp, nivel=1)
                elif tipo == 'solucion':
                    for paso in item.get('pasos', []):
                        agregar_texto_simple(text_frame, paso, nivel=1, color=COLORES['acento'])
                elif tipo == 'tabla':
                    agregar_tabla(
                        slide,
                        item.get('encabezados', []),
                        item.get('filas', []),
                        left=1, top=3, width=8, height=3
                    )


def agregar_diapositiva_contenido(prs, titulo, contenido_items):
    """
    Agrega una diapositiva con título y contenido.
    
    Args:
        prs: Objeto Presentation
        titulo (str): Título de la diapositiva
        contenido_items (list): Lista de items de contenido
    """
    slide_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)
    
    # Título
    titulo_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(1.2)
    )
    titulo_frame = titulo_box.text_frame
    p = titulo_frame.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORES['primario']
    p.alignment = PP_ALIGN.LEFT
    
    # Línea decorativa debajo del título
    line = slide.shapes.add_shape(
        1,  # Line shape
        Inches(0.5), Inches(1.4), Inches(9), Inches(0)
    )
    line.line.color.rgb = COLORES['acento']
    line.line.width = Pt(3)
    
    # Procesar contenido
    procesar_contenido(slide, contenido_items)


def generar_pptx(datos_yaml, output_path):
    """
    Genera un archivo PowerPoint desde datos YAML.
    
    Args:
        datos_yaml (dict): Datos del tema desde YAML
        output_path (str): Ruta donde guardar el archivo .pptx
    """
    # Crear presentación
    prs = crear_presentacion_base()
    
    # Agregar portada
    agregar_portada(
        prs,
        datos_yaml.get('tema', 'Sin título'),
        datos_yaml.get('subtitulo', '')
    )
    
    # Agregar diapositivas de contenido
    for diapositiva in datos_yaml.get('diapositivas', []):
        titulo = diapositiva.get('titulo', 'Sin título')
        contenido = diapositiva.get('contenido', [])
        agregar_diapositiva_contenido(prs, titulo, contenido)
    
    # Guardar presentación
    prs.save(output_path)
    print(f"✓ Presentación PowerPoint generada: {output_path}")


def procesar_tema(archivo_yaml, output_dir):
    """
    Procesa un tema completo: carga YAML y genera PowerPoint.
    
    Args:
        archivo_yaml (str): Ruta al archivo YAML del tema
        output_dir (str): Directorio base de salida
    """
    archivo_yaml = Path(archivo_yaml).resolve()
    print(f"\n→ Procesando: {archivo_yaml}")
    
    # Cargar datos
    datos = cargar_yaml(archivo_yaml)
    
    # Determinar directorio de salida manteniendo estructura
    try:
        directorio_base = Path(__file__).parent.resolve()
        clases_base = directorio_base / "clases"
        
        if archivo_yaml.is_relative_to(clases_base):
            relative_path = archivo_yaml.parent.relative_to(clases_base)
            pptx_output_dir = Path(output_dir) / relative_path
        else:
            pptx_output_dir = Path(output_dir)
    except (ValueError, AttributeError):
        pptx_output_dir = Path(output_dir)
    
    # Crear directorio de salida si no existe
    pptx_output_dir.mkdir(parents=True, exist_ok=True)
    
    # Generar nombre de salida
    nombre_base = archivo_yaml.stem
    output_pptx = pptx_output_dir / f"{nombre_base}.pptx"
    
    # Generar PowerPoint
    generar_pptx(datos, output_pptx)


def main():
    """Función principal del script."""
    parser = argparse.ArgumentParser(
        description='Generador de Presentaciones PowerPoint - Probabilidad y Estadística',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Ejemplos de uso:
  python generate_pptx.py                                    # Procesa TODOS los .yml en clases/**/
  python generate_pptx.py archivo.yml                        # Procesa un archivo específico
  python generate_pptx.py archivo1.yml archivo2.yml          # Procesa múltiples archivos
  python generate_pptx.py "clases/**/*.yml"                  # Procesa todos en clases (explícito)
  python generate_pptx.py "clases/probabilidad y estadistica/*.yml"  # Solo probabilidad
  python generate_pptx.py "clases/programacion_e_informatica/*.yml"  # Solo programación
  python generate_pptx.py -o pptx_output archivo.yml         # Especifica directorio de salida

Estructura de salida:
  - Archivos .pptx van a: pptx/[materia]/ (mantiene estructura de clases/)
        """
    )
    parser.add_argument(
        'archivos',
        nargs='*',
        help='Archivos YAML a procesar. Soporta patrones con * (ej: *.yml, clases/*.yml)'
    )
    parser.add_argument(
        '-o', '--output',
        dest='output_dir',
        default=None,
        help='Directorio de salida para los archivos generados (default: pptx/)'
    )
    
    args = parser.parse_args()
    
    # Configuración
    directorio_base = Path(__file__).parent
    output_dir = Path(args.output_dir) if args.output_dir else directorio_base / "pptx"
    clases_dir = directorio_base / "clases"
    
    # Determinar qué archivos procesar
    if args.archivos:
        # Expandir patrones con glob
        temas = []
        for patron in args.archivos:
            # Si el patrón contiene *, expandirlo
            if '*' in patron or '?' in patron:
                archivos_encontrados = glob.glob(patron, recursive=True)
                temas.extend([Path(f) for f in archivos_encontrados if f.endswith(('.yml', '.yaml'))])
            else:
                # Archivo directo
                temas.append(Path(patron))
    else:
        # Buscar todos los archivos .yml en clases/**/*
        temas = list(clases_dir.glob('**/*.yml'))
        if not temas:
            # Fallback a archivos por defecto de probabilidad y estadística
            prob_dir = clases_dir / "probabilidad y estadistica"
            temas = [
                prob_dir / "0-introduccion.yml",
                prob_dir / "1-tablas_graficos.yml",
                prob_dir / "2-medidas_posicion.yml",
                prob_dir / "3-reglas_probabilidades.yml"
            ]
    
    print("=" * 60)
    print("Generador de Presentaciones PowerPoint")
    print("Probabilidad y Estadística")
    print("=" * 60)
    
    # Verificar que hay archivos para procesar
    if not temas:
        print("✗ Error: No se encontraron archivos YAML para procesar")
        print("Usa --help para ver ejemplos de uso")
        return
    
    # Procesar cada tema
    procesados = 0
    errores = 0
    
    for tema in temas:
        if tema.exists():
            try:
                procesar_tema(tema, output_dir)
                procesados += 1
            except Exception as e:
                print(f"✗ Error procesando {tema.name}: {e}")
                errores += 1
        else:
            print(f"⚠ Advertencia: No se encontró {tema}")
    
    print("\n" + "=" * 60)
    print(f"✓ Proceso completado.")
    print(f"  {procesados} archivo(s) procesado(s) exitosamente.")
    if errores > 0:
        print(f"  {errores} archivo(s) con errores.")
    print(f"  Revisa el directorio: {output_dir}")
    print("=" * 60)


if __name__ == "__main__":
    main()
